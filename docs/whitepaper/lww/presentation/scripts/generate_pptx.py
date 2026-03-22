#!/usr/bin/env python3
"""
lww 발표용 PPTX 생성 스크립트
Usage: python generate_pptx.py --data ../content/slides-data.json --output ../slides/lww-pitch-deck.pptx [--with-images]
"""
import argparse
import json
import os
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx")
    raise

import re

def normalize_text(text: str) -> str:
    """한국어 텍스트 공백 정규화 (whitespace만, NFKC 제외)"""
    # 숫자+한글단위 사이 공백 제거: "3 분" -> "3분"
    text = re.sub(r'(\d)\s+(분|초|시간|일|주|개월|년|만|억|천|백|명|원|개|회|건|%)', r'\1\2', text)
    # 쉼표/마침표 앞 공백 제거: "면접 ," -> "면접,"
    text = re.sub(r'\s+([,.])', r'\1', text)
    # 연속 공백 정규화
    text = re.sub(r' {2,}', ' ', text)
    return text.strip()

# lww 브랜드 컬러
TEAL = RGBColor(0x0D, 0x94, 0x88)       # #0D9488
TEAL_LIGHT = RGBColor(0xF0, 0xFD, 0xFA)  # #F0FDFA
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x1A)       # #1A1A1A
GREY = RGBColor(0x6B, 0x72, 0x80)       # #6B7280

# 슬라이드 크기: 16:9 (1920x1080 기준 EMU)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def add_background(slide, color: RGBColor = None):
    """슬라이드 배경색 설정"""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color or WHITE


def add_textbox(slide, text, left, top, width, height,
                font_size=24, bold=False, color=None, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = normalize_text(text)
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color or DARK
    return txBox


def add_body_textbox(slide, lines, left, top, width, height, font_size=20):
    """바디 텍스트: 줄마다 paragraph 분리 + 간격 설정"""
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    from lxml import etree
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        # 줄 간격 및 아래 여백 설정
        p.space_after = Pt(10)
        p.line_spacing = Pt(font_size * 1.4)
        run = p.add_run()
        run.text = normalize_text(line)
        run.font.size = Pt(font_size)
        run.font.bold = False
        run.font.color.rgb = DARK
    return txBox


def add_cover_slide(prs, slide_data):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide, TEAL)

    # 메인 헤드라인
    add_textbox(
        slide, slide_data["headline"],
        Inches(1), Inches(2), Inches(11.3), Inches(2),
        font_size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER
    )
    # 서브헤드라인
    add_textbox(
        slide, slide_data["subheadline"],
        Inches(1), Inches(4), Inches(11.3), Inches(1),
        font_size=28, color=TEAL_LIGHT, align=PP_ALIGN.CENTER
    )
    # 바디
    add_textbox(
        slide, slide_data["body"],
        Inches(1), Inches(5.2), Inches(11.3), Inches(0.8),
        font_size=20, color=WHITE, align=PP_ALIGN.CENTER
    )
    return slide


def add_content_slide(prs, slide_data, with_images=False, images_dir=None, full_width=False):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide, WHITE)

    # 상단 Teal 바
    from pptx.util import Inches, Pt
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()

    # 전폭 모드: 텍스트 영역을 슬라이드 전체로 확장
    text_w = Inches(12.3) if full_width else Inches(7.5)
    body_w = Inches(12.3) if full_width else Inches(7.2)

    # 제목
    add_textbox(
        slide, slide_data["title"],
        Inches(0.5), Inches(0.25), Inches(12), Inches(0.5),
        font_size=14, color=TEAL
    )

    # 헤드라인
    add_textbox(
        slide, slide_data["headline"],
        Inches(0.5), Inches(0.9), text_w, Inches(1.2),
        font_size=28, bold=True, color=DARK
    )

    # 서브헤드라인
    add_textbox(
        slide, slide_data.get("subheadline", ""),
        Inches(0.5), Inches(2.1), text_w, Inches(0.9),
        font_size=16, color=GREY
    )

    # 바디 (멀티라인, bullet + 간격)
    raw_body = slide_data.get("body", "").replace("\\n", "\n")
    import re
    def needs_bullet(line):
        return not re.match(r'^[①②③④⑤⑥⑦⑧⑨\[\d]', line.strip())
    body_lines = [f"• {line}" if needs_bullet(line) else line for line in raw_body.split("\n") if line.strip()]
    add_body_textbox(
        slide, body_lines,
        Inches(0.5), Inches(3.2), body_w, Inches(3.2),
        font_size=20
    )

    # Phase 3: slide_id별 추가 도형 (이미지 없는 모드에서만)
    if not with_images:
        slide_id = slide_data["id"]
        if slide_id == 5:
            _add_funnel_shapes(slide)
        elif slide_id == 6:
            _add_matrix_shapes(slide)
        elif slide_id == 10:
            _add_timeline_shapes(slide)
        elif slide_id == 11:
            add_background(slide, TEAL_LIGHT)

    # 이미지 (옵션)
    if with_images and images_dir:
        img_filename = slide_data.get("image", "")
        img_path = Path(images_dir) / img_filename
        if img_path.exists():
            slide.shapes.add_picture(
                str(img_path),
                Inches(8), Inches(1),
                width=Inches(4.8)
            )
        else:
            # Placeholder 박스
            placeholder = slide.shapes.add_shape(
                1, Inches(8), Inches(1), Inches(4.8), Inches(5.5)
            )
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = TEAL_LIGHT
            placeholder.line.color.rgb = TEAL
            add_textbox(
                slide, f"[이미지]\n{img_filename}",
                Inches(8.1), Inches(3), Inches(4.6), Inches(1.5),
                font_size=14, color=TEAL, align=PP_ALIGN.CENTER
            )

    return slide


def add_cta_slide(prs, slide_data):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide, TEAL)

    add_textbox(
        slide, slide_data["headline"],
        Inches(1), Inches(1.5), Inches(11.3), Inches(1.5),
        font_size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER
    )
    add_textbox(
        slide, slide_data.get("subheadline", ""),
        Inches(1), Inches(3.2), Inches(11.3), Inches(1),
        font_size=24, color=TEAL_LIGHT, align=PP_ALIGN.CENTER
    )

    # URL 강조 박스
    url_text = slide_data.get("body", "")
    box_w = Inches(8)
    box_left = (SLIDE_WIDTH - box_w) // 2
    box_top = Inches(4.5)
    box_h = Inches(1.2)
    url_box = slide.shapes.add_shape(
        1, box_left, box_top, box_w, box_h
    )
    url_box.fill.solid()
    url_box.fill.fore_color.rgb = TEAL_LIGHT
    url_box.line.fill.background()

    add_textbox(
        slide, url_text,
        box_left + Inches(0.2), box_top + Inches(0.15),
        box_w - Inches(0.4), box_h - Inches(0.3),
        font_size=26, bold=True, color=TEAL, align=PP_ALIGN.CENTER
    )
    return slide


def _add_funnel_shapes(slide):
    """시장규모(slide 5): 3단 깔때기 사각형"""
    base_left = Inches(8.5)
    labels = [
        ("TAM  약 100만명", Inches(4.2), TEAL_LIGHT, TEAL),
        ("SAM  약 40만명",  Inches(3.0), TEAL_LIGHT, TEAL),
        ("SOM  25,000명",   Inches(1.8), TEAL,       WHITE),
    ]
    for i, (label, width, fill, text_color) in enumerate(labels):
        top = Inches(1.5) + Inches(1.7) * i
        left = base_left + (Inches(4.2) - width) // 2
        shape = slide.shapes.add_shape(1, left, top, width, Inches(1.3))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = TEAL
        from pptx.util import Pt as _Pt
        shape.line.width = _Pt(1)
        add_textbox(slide, label, left + Inches(0.1), top + Inches(0.25),
                    width - Inches(0.2), Inches(0.8),
                    font_size=14, bold=True, color=text_color, align=PP_ALIGN.CENTER)


def _add_matrix_shapes(slide):
    """경쟁사(slide 6): 2×2 포지셔닝 매트릭스"""
    cx, cy = Inches(10.7), Inches(3.9)
    half = Inches(2.2)
    from pptx.util import Pt as _Pt

    # 수평선
    h = slide.shapes.add_shape(1, cx - half, cy, half * 2, _Pt(2))
    h.fill.solid(); h.fill.fore_color.rgb = GREY; h.line.fill.background()
    # 수직선
    v = slide.shapes.add_shape(1, cx, cy - half, _Pt(2), half * 2)
    v.fill.solid(); v.fill.fore_color.rgb = GREY; v.line.fill.background()

    # 축 레이블
    add_textbox(slide, "양면 플랫폼 ↑",
                cx - half, cy - half - Inches(0.35),
                half * 2, Inches(0.3), font_size=10, color=GREY, align=PP_ALIGN.CENTER)
    add_textbox(slide, "AI 에이전트 →",
                cx + half + Inches(0.05), cy - Inches(0.15),
                Inches(1.3), Inches(0.3), font_size=10, color=GREY)

    # Fint 마커 (우상단)
    m = slide.shapes.add_shape(1, cx + Inches(0.5), cy - Inches(1.3), Inches(1.1), Inches(0.55))
    m.fill.solid(); m.fill.fore_color.rgb = TEAL; m.line.fill.background()
    add_textbox(slide, "Fint", cx + Inches(0.55), cy - Inches(1.25),
                Inches(1.0), Inches(0.45), font_size=16, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)


def _add_timeline_shapes(slide):
    """로드맵(slide 10): 수평 타임라인"""
    from pptx.util import Pt as _Pt
    line_y = Inches(3.6)
    x_start = Inches(8.2)
    line_w = Inches(4.6)

    # 수평선
    line = slide.shapes.add_shape(1, x_start, line_y, line_w, _Pt(3))
    line.fill.solid(); line.fill.fore_color.rgb = TEAL; line.line.fill.background()

    milestones = [
        ("MVP\n3/22",    0.0),
        ("Phase 1\n3/27", 0.45),
        ("Phase 2\n4/1",  0.9),
    ]
    dot = Inches(0.3)
    for label, ratio in milestones:
        mx = x_start + int(line_w * ratio)
        d = slide.shapes.add_shape(1, mx - dot // 2, line_y - dot // 2, dot, dot)
        d.fill.solid(); d.fill.fore_color.rgb = TEAL; d.line.fill.background()
        add_textbox(slide, label,
                    mx - Inches(0.55), line_y + Inches(0.25),
                    Inches(1.1), Inches(0.8),
                    font_size=11, color=DARK, align=PP_ALIGN.CENTER)


def main():
    parser = argparse.ArgumentParser(description="lww PPTX 생성")
    parser.add_argument("--data", required=True, help="slides-data.json 경로")
    parser.add_argument("--output", required=True, help="출력 PPTX 경로")
    parser.add_argument("--with-images", action="store_true", help="이미지 삽입")
    args = parser.parse_args()

    data_path = Path(args.data)
    output_path = Path(args.output)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    with open(data_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    images_dir = data_path.parent.parent / "images" if args.with_images else None

    prs = create_presentation()

    for slide_data in data["slides"]:
        slide_id = slide_data["id"]
        print(f"  슬라이드 {slide_id:02d}: {slide_data['title']}")

        if slide_id == 1:
            add_cover_slide(prs, slide_data)
        elif slide_id == 12:
            add_cta_slide(prs, slide_data)
        else:
            add_content_slide(prs, slide_data,
                              with_images=args.with_images,
                              images_dir=images_dir)

    prs.save(str(output_path))
    print(f"\n✅ PPTX 저장 완료: {output_path}")
    print(f"   슬라이드 수: {len(prs.slides)}")


if __name__ == "__main__":
    main()
